from pptx import Presentation
from datetime import datetime
import tempfile
import os

def processar_texto(texto_ou_arquivo):
    """
    Processa texto natural em formato simples:
    - Linha 1: Título da apresentação
    - Linha 2: Subtítulo  
    - Parágrafos seguintes (separados por linha vazia): Slides individuais
    """
    
    # Se for um arquivo, leia o conteúdo
    if isinstance(texto_ou_arquivo, str) and texto_ou_arquivo.endswith('.txt'):
        with open(texto_ou_arquivo, encoding='utf-8') as f:
            linhas = [linha.rstrip() for linha in f.readlines()]
    else:
        # Texto direto
        linhas = texto_ou_arquivo.strip().split('\n')
    
    # Remove linhas vazias do início e fim
    while linhas and not linhas[0].strip():
        linhas.pop(0)
    while linhas and not linhas[-1].strip():
        linhas.pop()
    
    if len(linhas) < 2:
        raise ValueError("Texto deve ter pelo menos título e subtítulo")
    
    # Primeira linha: título, Segunda linha: subtítulo
    titulo_apresentacao = linhas[0].strip()
    subtitulo = linhas[1].strip()
    
    # Gerar mês/ano atual automaticamente
    agora = datetime.now()
    meses = {
        1: 'Janeiro', 2: 'Fevereiro', 3: 'Março', 4: 'Abril', 
        5: 'Maio', 6: 'Junho', 7: 'Julho', 8: 'Agosto',
        9: 'Setembro', 10: 'Outubro', 11: 'Novembro', 12: 'Dezembro'
    }
    mes_ano = f"{meses[agora.month]}/{agora.year}"
    
    # Processar slides a partir da terceira linha
    slides_content = '\n'.join(linhas[2:]) if len(linhas) > 2 else ''
    
    # Separar parágrafos (cada parágrafo = um slide)
    paragrafos = []
    paragrafo_atual = []
    
    for linha in slides_content.split('\n'):
        if linha.strip():
            paragrafo_atual.append(linha.strip())
        else:
            if paragrafo_atual:
                paragrafos.append(paragrafo_atual)
                paragrafo_atual = []
    
    if paragrafo_atual:  # Adiciona último parágrafo
        paragrafos.append(paragrafo_atual)
    
    # Converter parágrafos em slides
    slides = []
    for i, paragrafo in enumerate(paragrafos, 1):
        if len(paragrafo) >= 1:
            titulo_slide = paragrafo[0]  # Primeira linha = título do slide
            conteudo_slide = '\n'.join(paragrafo[1:]) if len(paragrafo) > 1 else ''  # Resto = conteúdo
            
            slides.append({
                'NUMERO_TITULO': f'{i:02d}/{titulo_apresentacao}',
                'TITULO': titulo_slide,
                'CONTEUDO': conteudo_slide
            })
    
    return {
        'MES_ANO': mes_ano,
        'TITULO_APRESENTACAO': titulo_apresentacao,
        'SUBTITULO': subtitulo,
        'slides': slides
    }

def gerar_apresentacao_pptx(dados, modelo_path, saida_path):
    """
    Gera arquivo PowerPoint a partir dos dados processados
    """
    prs = Presentation(modelo_path)
    layout_capa = prs.slide_layouts[0]
    layout_conteudo = prs.slide_layouts[1]

    # Remove slides existentes do modelo
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    # Criar slide de capa
    slide_capa = prs.slides.add_slide(layout_capa)
    for shape in slide_capa.placeholders:
        ph_idx = shape.placeholder_format.idx
        
        # Preenchimento baseado nos índices específicos do modelo
        if ph_idx == 13:  # Título principal
            shape.text = dados.get('TITULO_APRESENTACAO', '')
        elif ph_idx == 12:  # Mês/Ano  
            shape.text = dados.get('MES_ANO', '')
        elif ph_idx == 15:  # Subtítulo
            shape.text = dados.get('SUBTITULO', '')

    # Criar slides de conteúdo
    for slide_info in dados['slides']:
        slide = prs.slides.add_slide(layout_conteudo)
        
        for shape in slide.placeholders:
            ph_idx = shape.placeholder_format.idx
            
            # Preenchimento dos placeholders de conteúdo
            if ph_idx == 13:  # Título do slide
                shape.text = slide_info.get('TITULO', '')
            elif ph_idx == 12:  # Número/Título  
                shape.text = slide_info.get('NUMERO_TITULO', '')
            elif ph_idx == 14:  # Conteúdo principal
                shape.text = slide_info.get('CONTEUDO', '')

    prs.save(saida_path)

def gerar_nome_arquivo(titulo):
    """
    Gera nome de arquivo limpo baseado no título da apresentação
    """
    titulo_limpo = ''.join(c for c in titulo if c.isalnum() or c in (' ', '-', '_')).strip()
    return f"{titulo_limpo.replace(' ', '_')}.pptx"

def criar_apresentacao(texto_entrada, modelo='modelo.pptx'):
    """
    Função principal: recebe texto e gera apresentação PowerPoint
    
    Args:
        texto_entrada: Texto natural ou caminho para arquivo .txt
        modelo: Caminho para arquivo modelo PowerPoint
        
    Returns:
        Nome do arquivo gerado
    """
    
    # Processar o texto
    dados = processar_texto(texto_entrada)
    
    # Gerar nome do arquivo baseado no título
    arquivo_saida = gerar_nome_arquivo(dados['TITULO_APRESENTACAO'])
    
    # Criar apresentação
    gerar_apresentacao_pptx(dados, modelo, arquivo_saida)
    
    print(f"✅ Apresentação '{dados['TITULO_APRESENTACAO']}' criada com {len(dados['slides'])} slides")
    print(f"📁 Arquivo salvo: {arquivo_saida}")
    
    return arquivo_saida

def processar_mensagem_teams(texto_usuario, modelo='modelo.pptx'):
    """
    Função específica para integração com Teams
    Recebe texto natural e retorna resultado estruturado
    """
    try:
        arquivo_gerado = criar_apresentacao(texto_usuario, modelo)
        
        dados = processar_texto(texto_usuario)
        return {
            'sucesso': True,
            'arquivo': arquivo_gerado,
            'titulo': dados['TITULO_APRESENTACAO'],
            'slides': len(dados['slides']),
            'mensagem': f"Apresentação '{dados['TITULO_APRESENTACAO']}' gerada com {len(dados['slides'])} slides"
        }
        
    except Exception as e:
        return {
            'sucesso': False,
            'erro': str(e),
            'mensagem': f"Erro ao gerar apresentação: {str(e)}"
        }

# Exemplo de uso
if __name__ == "__main__":
    
    # Exemplo com texto direto
    texto_exemplo = """Vendas Q1 2026
Análise de Performance

Resultados Globais
Ultrapassamos a meta trimestral em 15%, alcançando R$ 2.5M em vendas totais.
Teste 2

Performance por Região  
Norte: +20% | Sul: +12% | Sudeste: +18% | Nordeste: +8%

Próximos Passos
Focar expansão no Nordeste e manter momentum das demais regiões."""
    
    # Gera apresentação
    criar_apresentacao(texto_exemplo)